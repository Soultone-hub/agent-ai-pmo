import fitz  # PyMuPDF
from typing import cast as type_cast
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from openpyxl import load_workbook
from pptx import Presentation
from pptx.shapes.base import BaseShape
import os
import email
from email import policy


# ── Helper commun : rendu d'un tableau en Markdown ────────────────────────────
def _rows_to_markdown(rows: list) -> str:
    """Convertit une liste de lignes (chacune = liste de cellules) en tableau
    Markdown. Ignore les lignes entièrement vides et harmonise le nombre de
    colonnes pour produire un tableau GFM valide."""
    cleaned: list[list[str]] = []
    for row in rows or []:
        cells = [(str(c) if c is not None else "").replace("\n", " ").strip() for c in row]
        if any(cells):
            cleaned.append(cells)
    if not cleaned:
        return ""
    ncols = max(len(r) for r in cleaned)
    cleaned = [r + [""] * (ncols - len(r)) for r in cleaned]
    lines = [
        "| " + " | ".join(cleaned[0]) + " |",
        "| " + " | ".join(["---"] * ncols) + " |",
    ]
    for r in cleaned[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


# ── PDF ───────────────────────────────────────────────────────────────────────
def _is_real_table(rows: list) -> bool:
    """Distingue un vrai tableau de données d'un simple bloc de mise en page.
    Critères : ≥ 2 lignes non vides, ≥ 2 colonnes, et au moins la moitié des
    cellules remplies (les flyers/mises en page produisent des « tableaux »
    quasi vides qu'il ne faut pas rendre en Markdown)."""
    cleaned = [
        [(str(c) if c is not None else "").strip() for c in r]
        for r in (rows or [])
    ]
    cleaned = [r for r in cleaned if any(r)]
    if len(cleaned) < 2:
        return False
    ncols = max(len(r) for r in cleaned)
    if ncols < 2:
        return False
    total = sum(len(r) for r in cleaned)
    filled = sum(1 for r in cleaned for c in r if c)
    return total > 0 and (filled / total) >= 0.6


def _parse_pdf_fallback(file_path: str) -> str:
    """Extraction PDF native PyMuPDF : vrais tableaux convertis en Markdown +
    texte situé HORS de ces tableaux, le tout ordonné par position verticale et
    SANS duplication. Utilisé quand pymupdf4llm est indisponible (ex. Py 3.14)."""
    # Silence les avertissements MuPDF non bloquants (ex. « No common ancestor… »)
    try:
        fitz.TOOLS.mupdf_display_errors(False)
    except Exception:
        pass

    doc = fitz.open(file_path)
    parts: list[str] = []

    for page in doc:
        # 1. Détecter les tableaux, ne garder que les VRAIS tableaux de données
        try:
            found = page.find_tables()
            candidates = list(found.tables) if found else []
        except Exception:
            candidates = []

        real_tables: list[tuple[object, object, str]] = []  # (table, rect, markdown)
        for table in candidates:
            try:
                data = table.extract()
            except Exception:
                continue
            if _is_real_table(data):
                md = _rows_to_markdown(data)
                if md:
                    real_tables.append((table, fitz.Rect(table.bbox), md))
        table_rects = [rect for (_, rect, _) in real_tables]

        items: list[tuple[float, str]] = []  # (position_y, contenu)

        # 2. Blocs de texte situés HORS des vrais tableaux (évite la duplication)
        for block in page.get_text("blocks"):
            x0, y0, x1, y1, btext = block[0], block[1], block[2], block[3], block[4]
            btype = block[6] if len(block) > 6 else 0
            if btype != 0 or not str(btext).strip():
                continue
            bbox = fitz.Rect(x0, y0, x1, y1)
            area = bbox.get_area()
            inside_table = False
            for tr in table_rects:
                inter = bbox & tr
                if not inter.is_empty and area > 0 and (inter.get_area() / area) > 0.5:
                    inside_table = True
                    break
            if not inside_table:
                items.append((y0, str(btext).strip()))

        # 3. Tableaux rendus en Markdown
        for _, rect, md in real_tables:
            items.append((rect.y0, "\n" + md + "\n"))

        # 4. Restituer dans l'ordre vertical de la page
        items.sort(key=lambda it: it[0])
        parts.extend(content for _, content in items)

    return "\n\n".join(parts)


def parse_pdf(file_path: str) -> str:
    try:
        import pymupdf4llm
        # Markdown direct (préserve tableaux + structure) — actif sous Python 3.11
        md_text = str(pymupdf4llm.to_markdown(file_path)).strip()
        if md_text:
            return md_text
    except Exception:
        pass
    # Fallback natif PyMuPDF (Python 3.14, ou échec de pymupdf4llm)
    return _parse_pdf_fallback(file_path)


# ── DOCX ──────────────────────────────────────────────────────────────────────
def _iter_docx_blocks(parent):
    """Itère les paragraphes ET les tableaux d'un document Word dans leur ordre
    d'apparition (python-docx expose paragraphs et tables séparément, ce qui
    perd l'ordre — on parcourt donc directement le corps XML)."""
    body = parent.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _docx_table_to_markdown(table) -> str:
    rows = [[cell.text for cell in row.cells] for row in table.rows]
    return _rows_to_markdown(rows)


def parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    parts: list[str] = []
    for block in _iter_docx_blocks(doc):
        if isinstance(block, Paragraph):
            txt = block.text.strip()
            if txt:
                parts.append(txt)
        else:  # Table
            md = _docx_table_to_markdown(block)
            if md:
                parts.append("\n" + md + "\n")
    return "\n".join(parts)

def parse_xlsx(file_path: str) -> str:
    wb = load_workbook(file_path)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text += f"Feuille: {sheet}\n"
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                text += row_text + "\n"
    return text

def parse_pptx(file_path: str) -> str:
    """Extrait le texte de toutes les diapositives d'un fichier PPTX."""
    prs = Presentation(file_path)
    text = ""
    for slide_num, slide in enumerate(prs.slides, start=1):
        text += f"\n--- Diapositive {slide_num} ---\n"
        for shape in slide.shapes:
            # Titre et contenus texte
            if shape.has_text_frame:
                from pptx.shapes.autoshape import Shape as AutoShape
                tf_shape = type_cast(AutoShape, shape)
                for para in tf_shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        text += line + "\n"
            # Tables
            if shape.has_table:
                from pptx.shapes.graphfrm import GraphicFrame
                tbl_shape = type_cast(GraphicFrame, shape)
                for row in tbl_shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        text += row_text + "\n"
        # Notes du présentateur
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide is not None and notes_slide.notes_text_frame is not None:
                notes = notes_slide.notes_text_frame.text.strip()
                if notes:
                    text += f"[Notes] {notes}\n"
    return text


def parse_document(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".xlsx":
        return parse_xlsx(file_path)
    elif ext == ".eml":
        return parse_eml(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".ppt":
        raise ValueError("Le format .ppt (PowerPoint ancien) n'est pas supporté. Convertissez en .pptx.")
    else:
        raise ValueError(f"Format non supporté: {ext}")
    

def parse_eml(file_path: str) -> str:
    with open(file_path, 'rb') as f:
        msg = email.message_from_binary_file(f, policy=policy.default)
    
    text = f"De: {msg['from']}\n"
    text += f"À: {msg['to']}\n"
    text += f"Sujet: {msg['subject']}\n"
    text += f"Date: {msg['date']}\n\n"
    
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                text += part.get_content()
    else:
        text += msg.get_content()
    
    return text